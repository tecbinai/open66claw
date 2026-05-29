# -*- coding: utf-8 -*-
"""
woclaw PPT v5 — 顶级重设计
Page 1: woclaw 是什么 — 三卡片，要点两行布局（修复空白）
Page 2: 为什么是 woclaw — 左侧四象限定位图 + 右侧核心优势
        四象限：X=开放性，Y=安全稳定性
        三产品：大厂Claw(左上) / OpenClaw(右下) / woclaw(右上最优)
        无红色，底部内容不遮盖页脚
"""
import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from lxml import etree
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = r'D:\newopenclaw\调研\《英迈AI生态引擎：从异构硬件到智能体经济的全栈实践》6-20.pptx'
OUT = r'D:\newopenclaw\调研\woclaw-slides-v5.pptx'

W, H = 13.33, 7.50
E = 914400

def i(x):  return int(x * E)
def p(x):  return int(x * 12700)

# ── 颜色 ────────────────────────────────────────────────────────
BLU  = RGBColor(0x00, 0x77, 0xD4)  # 英迈蓝
DBLU = RGBColor(0x00, 0x4F, 0x9A)  # 深蓝
ABLU = RGBColor(0x00, 0xA1, 0xE4)  # 亮蓝
MBLU = RGBColor(0x22, 0x76, 0xBB)  # 中蓝
BLK  = RGBColor(0x1A, 0x1A, 0x2E)  # 近黑
GRY  = RGBColor(0x4A, 0x55, 0x68)  # 中灰
LGRY = RGBColor(0xF3, 0xF5, 0xF8)  # 浅灰背景
MGRY = RGBColor(0xB0, 0xBE, 0xCC)  # 竞品气泡灰
BRD  = RGBColor(0xCC, 0xD5, 0xE0)  # 边框
WHT  = RGBColor(0xFF, 0xFF, 0xFF)
ORG  = RGBColor(0xF2, 0x6B, 0x43)  # 英迈橙


def get_layout(prs):
    for l in prs.slide_layouts:
        if "Global" in l.name: return l
    for l in prs.slide_layouts:
        if "标准" in l.name: return l
    return prs.slide_layouts[0]


def nuke_ph(slide):
    NS = 'http://schemas.openxmlformats.org/presentationml/2006/main'
    spTree = slide.shapes._spTree
    kill = []
    for sp in spTree.findall(f'.//{{{NS}}}sp'):
        nvPr = sp.find(f'.//{{{NS}}}nvPr')
        if nvPr is not None:
            ph = nvPr.find(f'{{{NS}}}ph')
            if ph is not None and ph.get('type', '') not in ('sldNum', 'dt', 'ftr'):
                kill.append(sp)
    for sp in kill:
        try: spTree.remove(sp)
        except: pass


def r(slide, l, t, w, h, fill=None, ln=None, lw=0.5):
    s = slide.shapes.add_shape(1, i(l), i(t), i(w), i(h))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if ln: s.line.color.rgb = ln; s.line.width = p(lw)
    return s


def rr(slide, l, t, w, h, fill=None, ln=None, lw=0.5, adj=3500):
    s = slide.shapes.add_shape(5, i(l), i(t), i(w), i(h))
    ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
    pg = s._element.spPr.find(f'{{{ns}}}prstGeom')
    if pg is not None:
        al = pg.find(f'{{{ns}}}avLst')
        if al is not None:
            for g in al.findall(f'{{{ns}}}gd'): al.remove(g)
            gd = etree.SubElement(al, f'{{{ns}}}gd')
            gd.set('name', 'adj'); gd.set('fmla', f'val {adj}')
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if ln: s.line.color.rgb = ln; s.line.width = p(lw)
    return s


def ellipse(slide, cx, cy, rw, rh, fill=None, ln=None, lw=0.8):
    s = slide.shapes.add_shape(9, i(cx - rw), i(cy - rh), i(rw*2), i(rh*2))
    s.fill.solid() if fill else s.fill.background()
    if fill: s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if ln: s.line.color.rgb = ln; s.line.width = p(lw)
    return s


def tb(slide, text, l, t, w, h, sz=16, bold=False, clr=GRY,
       align=PP_ALIGN.LEFT, wrap=True):
    b = slide.shapes.add_textbox(i(l), i(t), i(w), i(h))
    b.word_wrap = wrap
    tf = b.text_frame; tf.word_wrap = wrap
    pg = tf.paragraphs[0]; pg.alignment = align
    rv = pg.add_run(); rv.text = text
    rv.font.size = Pt(sz); rv.font.bold = bold; rv.font.color.rgb = clr
    return b, tf


def mx(slide, runs, l, t, w, h, align=PP_ALIGN.LEFT, wrap=True):
    b = slide.shapes.add_textbox(i(l), i(t), i(w), i(h))
    b.word_wrap = wrap
    tf = b.text_frame; tf.word_wrap = wrap
    pg = tf.paragraphs[0]; pg.alignment = align
    for text, sz, bold, clr in runs:
        rv = pg.add_run(); rv.text = text
        rv.font.size = Pt(sz); rv.font.bold = bold; rv.font.color.rgb = clr
    return b


def add_top_bar(slide):
    r(slide, 0, 0, W, 0.055, fill=BLU)


# ═══════════════════════════════════════════════════════════════════
# PAGE 1：woclaw 是什么
# 三卡片：安全合规 / 国产适配 / 开放扩展
# 要点采用两行布局（关键词行 + 描述行），彻底消除空白感
# ═══════════════════════════════════════════════════════════════════
def page1(prs):
    slide = prs.slides.add_slide(get_layout(prs))
    nuke_ph(slide)
    add_top_bar(slide)

    mx(slide, [
        ("woclaw  ", 28, True, BLU),
        ("是什么", 28, True, BLK),
    ], l=0.60, t=0.15, w=10.0, h=0.62)

    tb(slide, "企业 AI Agent 本地部署平台  ·  英迈AI生态的智能体落地工具",
       l=0.60, t=0.80, w=12.0, h=0.36, sz=14, clr=GRY)

    CW, CH = 3.96, 4.62
    CT = 1.26
    COLS = [0.38, 4.56, 8.74]
    TOP_COLORS = [BLU, ABLU, DBLU]

    cards = [
        {
            "big_num": "12/12",
            "big_sub": "工信部合规全覆盖",
            "title":   "安全合规",
            "points": [
                ("数据本地",  "100% 不出境，不上任何云"),
                ("全程可查",  "操作日志本地存储可追溯"),
                ("人工确认",  "危险操作必须人工二次确认"),
            ],
        },
        {
            "big_num": "15+",
            "big_sub": "国产大模型随时切换",
            "title":   "国产适配",
            "points": [
                ("主流模型",  "Kimi · 通义 · 豆包 · 硅基等"),
                ("全平台 IM", "飞书/钉钉/企微/微信/QQ"),
                ("中文界面",  "全中文，无需培训即可上手"),
            ],
        },
        {
            "big_num": "13000+",
            "big_sub": "技能插件生态",
            "title":   "开放扩展",
            "points": [
                ("零绑定",   "模型与 IM 随时可替换"),
                ("即装即用",  "ClawHub 技能一键安装"),
                ("周级交付",  "需求定制，快速响应"),
            ],
        },
    ]

    for ci, (cl, tc, card) in enumerate(zip(COLS, TOP_COLORS, cards)):
        rr(slide, cl, CT, CW, CH, fill=LGRY, ln=BRD, lw=0.4)
        r(slide, cl, CT, CW, 0.07, fill=tc)

        # 超大数字
        tb(slide, card["big_num"],
           l=cl+0.22, t=CT+0.14, w=CW-0.35, h=0.86,
           sz=38, bold=True, clr=tc)

        # 数字副标
        tb(slide, card["big_sub"],
           l=cl+0.22, t=CT+0.98, w=CW-0.35, h=0.30,
           sz=12, clr=GRY)

        # 分隔线
        r(slide, cl+0.22, CT+1.33, CW-0.44, 0.016, fill=BRD)

        # 卡片标题
        tb(slide, card["title"],
           l=cl+0.22, t=CT+1.38, w=CW-0.35, h=0.44,
           sz=20, bold=True, clr=BLK)

        # 三条要点（两行：关键词 + 描述）
        for ri, (kw, desc) in enumerate(card["points"]):
            py = CT + 1.94 + ri * 0.84
            r(slide, cl+0.22, py+0.06, 0.04, 0.62, fill=tc)
            tb(slide, kw,
               l=cl+0.34, t=py, w=CW-0.52, h=0.34,
               sz=15, bold=True, clr=BLK, wrap=False)
            tb(slide, desc,
               l=cl+0.34, t=py+0.33, w=CW-0.52, h=0.38,
               sz=12, clr=GRY, wrap=True)

    # 底部钉子句
    BT = CT + CH + 0.08
    r(slide, 0, BT, W, H - BT - 0.04, fill=DBLU)
    tb(slide, "本地部署的安全  ·  丰富的技能生态  ·  不被任何一家锁死",
       l=0.50, t=BT+0.12, w=W-1.0, h=0.46,
       sz=18, bold=True, clr=WHT, align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════════
# PAGE 2：为什么是 woclaw
#
# 左侧（6.50"）：四象限定位图
#   X轴：开放性（低→高）
#   Y轴：安全稳定性（低→高）
#   大厂Claw  → 左上（安全高 + 开放低）
#   OpenClaw  → 右下（开放高 + 安全低）
#   woclaw    → 右上（安全高 + 开放高）★最优
#   左下象限空置（标注"不存在"）
#
# 右侧（6.20"）：5条核心优势（带序号色块）
# 底部钉子句 y=6.08，不遮盖页脚
# ═══════════════════════════════════════════════════════════════════
def page2(prs):
    slide = prs.slides.add_slide(get_layout(prs))
    nuke_ph(slide)
    add_top_bar(slide)

    # 标题
    mx(slide, [
        ("为什么是 ", 28, True, BLK),
        ("woclaw", 28, True, BLU),
        (" ？", 28, True, BLK),
    ], l=0.60, t=0.15, w=10.0, h=0.62)

    tb(slide, "安全稳定  ·  开放灵活  ·  中国企业的唯一最优解",
       l=0.60, t=0.80, w=12.0, h=0.36, sz=14, clr=GRY)

    # ══════════════════════════════════════════════════════════════
    # 左侧：四象限图
    # 图区：x=0.28, y=1.16, w=6.50, h=4.76  底=5.92
    # ══════════════════════════════════════════════════════════════
    GX, GY, GW, GH = 0.28, 1.16, 6.50, 4.76
    GB = GY + GH  # 图底 y = 5.92

    # 图区白色背景
    rr(slide, GX, GY, GW, GH, fill=WHT, ln=BRD, lw=0.5)

    # 原点（坐标轴交叉点，图区正中）
    OX = GX + GW / 2   # 3.53
    OY = GY + GH / 2   # 3.54

    # ── 四象限背景色（极淡）──────────────────────────────────────
    # 右上象限（woclaw区）：极淡蓝
    r(slide, OX, GY, GX+GW-OX, OY-GY,
      fill=RGBColor(0xE8, 0xF4, 0xFF), ln=None)
    # 左上象限：极淡灰
    r(slide, GX, GY, OX-GX, OY-GY,
      fill=RGBColor(0xF5, 0xF7, 0xFA), ln=None)
    # 右下象限：极淡灰
    r(slide, OX, OY, GX+GW-OX, GB-OY,
      fill=RGBColor(0xF5, 0xF7, 0xFA), ln=None)
    # 左下象限：更淡灰（空置）
    r(slide, GX, OY, OX-GX, GB-OY,
      fill=RGBColor(0xF9, 0xFA, 0xFB), ln=None)

    # ── 坐标轴 ───────────────────────────────────────────────────
    # X轴（横线）
    r(slide, GX+0.10, OY-0.010, GW-0.20, 0.020, fill=RGBColor(0x88, 0xA0, 0xBB))
    # Y轴（竖线）
    r(slide, OX-0.010, GY+0.10, 0.020, GH-0.20, fill=RGBColor(0x88, 0xA0, 0xBB))

    # 轴标签
    # X轴两端
    tb(slide, "开放性低", l=GX+0.12, t=OY+0.04, w=1.10, h=0.26,
       sz=9, clr=RGBColor(0x88, 0xA0, 0xBB))
    tb(slide, "开放性高", l=GX+GW-1.22, t=OY+0.04, w=1.10, h=0.26,
       sz=9, clr=RGBColor(0x88, 0xA0, 0xBB), align=PP_ALIGN.RIGHT)
    # Y轴两端
    tb(slide, "安全性高", l=OX+0.06, t=GY+0.10, w=1.10, h=0.26,
       sz=9, clr=RGBColor(0x88, 0xA0, 0xBB))
    tb(slide, "安全性低", l=OX+0.06, t=GB-0.36, w=1.10, h=0.26,
       sz=9, clr=RGBColor(0x88, 0xA0, 0xBB))

    # ── 象限标签（角落，淡色）────────────────────────────────────
    # 右上
    tb(slide, "自主可控", l=OX+0.14, t=GY+0.14, w=1.60, h=0.30,
       sz=10, bold=True, clr=RGBColor(0x00, 0x77, 0xD4),
       align=PP_ALIGN.LEFT)
    # 左上
    tb(slide, "安全孤岛", l=GX+0.14, t=GY+0.14, w=1.60, h=0.30,
       sz=10, bold=False, clr=RGBColor(0xA0, 0xB0, 0xC0))
    # 右下
    tb(slide, "开放有险", l=OX+0.14, t=GB-0.44, w=1.60, h=0.30,
       sz=10, bold=False, clr=RGBColor(0xA0, 0xB0, 0xC0))
    # 左下（空置）
    tb(slide, "无意义区", l=GX+0.14, t=GB-0.44, w=1.60, h=0.30,
       sz=9, bold=False, clr=RGBColor(0xCC, 0xD5, 0xE0))

    # ── 产品气泡 ─────────────────────────────────────────────────
    # 大厂Claw：左上（开放低+安全高）→ (-0.52, +0.52) 偏移
    DC_X = OX - 0.52 * (OX - GX - 0.30)
    DC_Y = OY - 0.52 * (OY - GY - 0.30)
    ellipse(slide, DC_X, DC_Y, 0.52, 0.30,
            fill=RGBColor(0xD4, 0xDF, 0xEC), ln=MBLU, lw=0.6)
    tb(slide, "大厂 Claw",
       l=DC_X-0.52, t=DC_Y-0.18, w=1.04, h=0.24,
       sz=10, bold=True, clr=DBLU, align=PP_ALIGN.CENTER)
    tb(slide, "安全但封闭绑定",
       l=DC_X-0.60, t=DC_Y+0.10, w=1.20, h=0.22,
       sz=8, clr=GRY, align=PP_ALIGN.CENTER)

    # OpenClaw：右下（开放高+安全低）→ (+0.52, -0.52)
    OC_X = OX + 0.52 * (GX + GW - 0.30 - OX)
    OC_Y = OY + 0.52 * (GB - 0.30 - OY)
    ellipse(slide, OC_X, OC_Y, 0.52, 0.30,
            fill=RGBColor(0xD4, 0xDF, 0xEC), ln=MBLU, lw=0.6)
    tb(slide, "OpenClaw",
       l=OC_X-0.52, t=OC_Y-0.18, w=1.04, h=0.24,
       sz=10, bold=True, clr=DBLU, align=PP_ALIGN.CENTER)
    tb(slide, "开放但维护弱",
       l=OC_X-0.60, t=OC_Y+0.10, w=1.20, h=0.22,
       sz=8, clr=GRY, align=PP_ALIGN.CENTER)

    # woclaw：右上（开放高+安全高）→ 大蓝圆，最醒目
    WC_X = OX + 0.58 * (GX + GW - 0.30 - OX)
    WC_Y = OY - 0.58 * (OY - GY - 0.30)
    ellipse(slide, WC_X, WC_Y, 0.62, 0.36,
            fill=BLU, ln=DBLU, lw=1.0)
    tb(slide, "woclaw",
       l=WC_X-0.62, t=WC_Y-0.20, w=1.24, h=0.26,
       sz=12, bold=True, clr=WHT, align=PP_ALIGN.CENTER)
    tb(slide, "唯一最优解",
       l=WC_X-0.62, t=WC_Y+0.06, w=1.24, h=0.22,
       sz=9, bold=False, clr=RGBColor(0xB8, 0xD8, 0xFF),
       align=PP_ALIGN.CENTER)

    # ── 图区标题 ─────────────────────────────────────────────────
    tb(slide, "产品定位矩阵",
       l=GX+0.18, t=GY+0.12, w=2.20, h=0.28,
       sz=11, bold=True, clr=DBLU)

    # ══════════════════════════════════════════════════════════════
    # 右侧：核心优势列表（5条）
    # x=6.98, y=1.16, w=6.12
    # 每条高 0.82，间距 0.09，5条总高 = 5*0.82+4*0.09 = 4.46
    # 底部 y = 1.16+0.36+4.46 = 5.98  安全
    # ══════════════════════════════════════════════════════════════
    RX = 6.98
    RW = 6.12
    RT = 1.16

    tb(slide, "woclaw 核心优势",
       l=RX, t=RT+0.06, w=RW, h=0.34,
       sz=14, bold=True, clr=DBLU)

    advantages = [
        (BLU,  "数据100%本地",
                "完全自主部署，数据不出境、不上云，工信部合规12/12全覆盖"),
        (ABLU, "专业团队保障",
                "Woclaw团队7x24支持，持续迭代升级，企业级SLA保障"),
        (MBLU, "15+ 模型零绑定",
                "国产主流大模型随时切换，不被任何单一厂商锁死"),
        (DBLU, "全平台IM打通",
                "飞书/钉钉/企微/微信/QQ，一套部署全部贯通"),
        (BLU,  "13000+ 开放生态",
                "ClawHub插件市场，需求定制周级响应，无限扩展"),
    ]

    ADV_H = 0.82
    ADV_G = 0.09
    for ai, (nc, title, desc) in enumerate(advantages):
        ay = RT + 0.52 + ai * (ADV_H + ADV_G)

        # 整条背景
        rr(slide, RX, ay, RW, ADV_H, fill=LGRY, ln=BRD, lw=0.3)

        # 左侧色块 + 序号
        r(slide, RX, ay, 0.46, ADV_H, fill=nc)
        tb(slide, f"0{ai+1}",
           l=RX+0.02, t=ay+(ADV_H-0.36)/2, w=0.42, h=0.36,
           sz=14, bold=True, clr=WHT, align=PP_ALIGN.CENTER)

        # 强调竖线
        r(slide, RX+0.46, ay+0.10, 0.032, ADV_H-0.20, fill=nc)

        # 标题（粗体深色）
        tb(slide, title,
           l=RX+0.56, t=ay+0.08, w=RW-0.64, h=0.30,
           sz=14, bold=True, clr=BLK, wrap=False)

        # 描述（灰色小字）
        tb(slide, desc,
           l=RX+0.56, t=ay+0.38, w=RW-0.64, h=0.36,
           sz=11, clr=GRY, wrap=True)

    # ── 底部钉子句（y=6.08，不遮盖 Ingram 页脚）────────────────
    NT = 6.08
    NH = 0.58
    r(slide, 0, NT, W, NH, fill=DBLU)
    tb(slide,
       "数据安全自主  ·  模型自由切换  ·  专业团队保障  ·  这就是 woclaw",
       l=0.50, t=NT+0.07, w=W-1.0, h=0.44,
       sz=16, bold=True, clr=WHT, align=PP_ALIGN.CENTER)

    return slide


# ═══════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    solo = Presentation(SRC)
    page1(solo)
    page2(solo)
    sIdLst = solo.slides._sldIdLst
    for el in list(sIdLst)[:19]:
        sIdLst.remove(el)
    solo.save(OUT)
    print(f"Done: {OUT}")
